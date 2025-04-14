from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Function to set smaller font size for text frames
def set_smaller_font(text_frame, font_size=Pt(20)):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = font_size

# Create a presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
title = slide_1.shapes.title
subtitle = slide_1.placeholders[1]
title.text = "Analysis of Blocked Payments in Stripe"
subtitle.text = "Behavior of prepaid card payment blocks\n[Current Date]"

# Slide 2: Context
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout
title = slide_2.shapes.title
content = slide_2.shapes.placeholders[1].text_frame
title.text = "Context of the Analysis"
content.text = (
    "- Stripe blocks subscription payments made with prepaid cards.\n"
    "- Objective: Analyze the behavior of these blocks to make informed decisions.\n"
    "- Key Metrics: Blocked payments, blocked amounts, blocked users, resolved payments, resolved users, new blocked users."
)
set_smaller_font(content)  # Set smaller font size

# Slide 3: Key Metrics Definitions (Part 1)
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_3.shapes.title
content = slide_3.shapes.placeholders[1].text_frame
title.text = "Key Metrics Definitions (Part 1)"
content.text = (
    "- **real_block_payments:** Number of unique payment errors (not repeated attempts). "
    "For example, if a user tries to purchase a subscription 5 times and all attempts fail, "
    "it counts as 1 payment error, not 5.\n\n"
    "- **real_amount_blocked:** Total amount blocked due to payment errors. "
    "This follows the same logic as real_block_payments but focuses on the monetary value.\n\n"
    "- **blocked_users:** Number of unique users blocked. "
    "This metric counts each user only once, regardless of how many payment attempts they made.\n\n"
    "- **blocked_payments_resolved:** Number of blocked payments that were successfully resolved. "
    "This indicates how many payment errors were fixed after being blocked."
)
set_smaller_font(content)  # Set smaller font size

# Slide 4: Key Metrics Definitions (Part 2)
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_4.shapes.title
content = slide_4.shapes.placeholders[1].text_frame
title.text = "Key Metrics Definitions (Part 2)"
content.text = (
    "- **amount_resolved:** Total amount unblocked after resolving payment errors. "
    "This represents the monetary value of resolved payments.\n\n"
    "- **blocked_users_resolved:** Number of blocked users that were resolved. "
    "This indicates how many users were unblocked after their payment issues were fixed.\n\n"
    "- **new_users_blocked:** Number of new users blocked. "
    "A new user is defined as someone whose email is not registered in the database. "
    "If the email exists, we check if there are no successful purchase orders at least one day before the payment error. "
    "If there are no orders, the user is considered new.\n\n"
    "- **percentage_new_users_blocked (%):** Percentage of new blocked users relative to the total of new users and blocked users. "
    "It is calculated as: (new_users_blocked) / (total_new_users + blocked_users - blocked_users_resolved)."
)
set_smaller_font(content)  # Set smaller font size

# Slide 5: Total Data Summary
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_5.shapes.title
content = slide_5.shapes.placeholders[1].text_frame
title.text = "Summary of Total Data"
content.text = (
    "- Total Blocked Payments: 45\n"
    "- Total Amount Blocked: $843.65\n"
    "- Total Blocked Users: 38\n"
    "- Total Resolved Payments: 14\n"
    "- Total Amount Resolved: $272.53\n"
    "- Total Resolved Users: 14\n"
    "- Total New Blocked Users: 27\n"
    "- Percentage of New Blocked Users: 7.05%"
)
set_smaller_font(content)  # Set smaller font size

# Slide 6: Behavior Charts (Placeholder for Charts)
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_6.shapes.title
content = slide_6.shapes.placeholders[1].text_frame
title.text = "Blocked Payments Behavior"
content.text = (
    "Charts to be added:\n"
    "- Blocked payments per day.\n"
    "- Blocked amount per day.\n"
    "- Blocked users per day.\n"
    "- Resolved payments per day.\n"
    "- New blocked users per day."
)
set_smaller_font(content)  # Set smaller font size

# Slide 7: New vs Existing Users Analysis
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_7.shapes.title
content = slide_7.shapes.placeholders[1].text_frame
title.text = "New vs Existing Blocked Users"
content.text = (
    "- Total Blocked Users: 38\n"
    "- New Blocked Users: 27\n"
    "- Existing Blocked Users: 11 (38 - 27)\n"
    "- Conclusion: 71.05% of blocked users are new, while 28.95% are existing users."
)
set_smaller_font(content)  # Set smaller font size

# Slide 8: Conclusions and Recommendations
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_8.shapes.title
content = slide_8.shapes.placeholders[1].text_frame
title.text = "Conclusions and Recommendations"
content.text = (
    "- Resolved Payments: Out of 45 blocked payments, 14 were resolved (31.11% resolution rate).\n"
    "- Existing Users Blocked: 11 existing users were blocked, representing 28.95% of total blocked users.\n"
    "- Recommendations:\n"
    "  - Review blocking policies to avoid affecting existing users.\n"
    "  - Implement better mechanisms to identify new users and prevent unnecessary blocks.\n"
    "  - Continuously monitor the percentage of new blocked users to adjust strategies."
)
set_smaller_font(content)  # Set smaller font size

# Slide 9: Questions and Discussion
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title = slide_9.shapes.title
content = slide_9.shapes.placeholders[1].text_frame
title.text = "Questions and Discussion"
content.text = "Space for questions and discussion about the results and recommendations."
set_smaller_font(content)  # Set smaller font size

# Save the presentation
prs.save("Blocked_Payments_Analysis.pptx")

print("Presentation created successfully!")